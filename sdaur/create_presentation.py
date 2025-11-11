from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Define slide layout (blank)
blank_slide_layout = prs.slide_layouts[6]

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(14), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle_frame.paragraphs[0].font.size = Pt(28)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

    return slide

def add_question_slide(prs, question, image_path):
    """Add a slide with question as title and graph as content"""
    slide = prs.slides.add_slide(blank_slide_layout)

    # Question (title)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = question
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # Image (graph)
    if os.path.exists(image_path):
        # Center the image
        left = Inches(2)
        top = Inches(1.5)
        height = Inches(6.5)
        slide.shapes.add_picture(image_path, left, top, height=height)
    else:
        # Placeholder text if image doesn't exist
        placeholder = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(12), Inches(2))
        placeholder.text_frame.text = f"[Image: {os.path.basename(image_path)}]"
        placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        placeholder.text_frame.paragraphs[0].font.size = Pt(24)
        placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)

    return slide

def add_conclusion_slide(prs, text):
    """Add conclusion slide"""
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Conclusion"
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True

    # Content
    content_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(12), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = text
    content_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    content_frame.paragraphs[0].font.size = Pt(24)
    content_frame.paragraphs[0].line_spacing = 1.5

    return slide

# Title Slide
add_title_slide(prs,
                "Marine Traffic Data Analysis",
                "An Exploratory Study of Global Vessel Patterns")

# Question Slides (8 questions max)
plots_dir = "plots"

questions_and_plots = [
    ("What does global maritime traffic look like?",
     "05_geographic_distribution.png"),

    ("Do different vessel types have characteristic shapes?",
     "03_lw_ratio_by_shiptype.png"),

    ("How does cargo capacity relate to vessel size?",
     "04_dwt_vs_length.png"),

    ("How does speed vary geographically?",
     "06_speed_by_location.png"),

    ("Can heading-course differences identify vessels in distress?",
     "02_hc_diff_anomalies.png"),

    ("Which features show strong correlations?",
     "07_correlation_matrix.png"),

    ("What's missing in our data?",
     "08_missing_data_analysis.png"),

    ("Can we model dimensional relationships?",
     "12_dimensions_with_speed.png")
]

for question, plot_file in questions_and_plots:
    image_path = os.path.join(plots_dir, plot_file)
    add_question_slide(prs, question, image_path)

# Conclusion Slide
conclusion_text = """We reverse-engineered the MarineTraffic API and scraped 10,379 vessels globally.

Through exploratory analysis, we discovered systematic relationships between vessel
characteristics, identified operational patterns, and built predictive models.

Simple linear regression explained 51.8% of dimensional variance—satisfying academic
requirements while revealing the need for advanced machine learning approaches.

This static snapshot establishes the foundation for real-time maritime surveillance,
anomaly detection, and predictive analytics."""

add_conclusion_slide(prs, conclusion_text)

# Save presentation
output_path = "marine_traffic_presentation.pptx"
prs.save(output_path)
print(f"Presentation created: {output_path}")
